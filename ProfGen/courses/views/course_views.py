from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib.auth import login
from django.contrib.auth.views import LoginView
from ..models import Course, UploadedMaterial
from ..forms import CourseForm, UploadedMaterialForm
from django.http import HttpResponseForbidden
from django.core.exceptions import PermissionDenied
import os
from PyPDF2 import PdfReader
from docx import Document
from django.http import HttpResponseBadRequest
from pptx import Presentation

@login_required
def course_list(request):
    if not request.user.is_authenticated:
        return redirect('login')
    courses = Course.objects.filter(professor=request.user).all()
    return render(request, 'courses/course_list.html', {'courses': courses})

@login_required
def course_detail(request, course_id):
    if not request.user.is_authenticated:
        return redirect('login')
    course = get_object_or_404(Course, id=course_id)
    course_materials = UploadedMaterial.objects.filter(course=course)
    return render(request, 'courses/course_detail.html', {'course': course, 'course_materials': course_materials})

@login_required
def create_course(request):
    if request.method == 'POST':
        form = CourseForm(request.POST)
        if form.is_valid():
            course = form.save(commit=False)  # Do not immediately save to DB
            course.professor = request.user  # Set the professor to the current user
            course.save()  # Now save the course
            return redirect('course_detail', course_id=course.id)
    else:
        form = CourseForm()
    return render(request, 'courses/create_course.html', {'form': form})


@login_required
def edit_course(request, course_id):
    course = get_object_or_404(Course, id=course_id)
    if request.method == 'POST':
        form = CourseForm(request.POST, instance=course)
        if form.is_valid():
            course = form.save()
            return redirect('course_detail', course_id=course.id)
    else:
        form = CourseForm(instance=course)
    return render(request, 'courses/edit_course.html', {'form': form, 'course': course})

# @login_required
# def upload_material(request, course_id):
#     course = get_object_or_404(Course, id=course_id, professor=request.user)
#     if request.method == 'POST':
#         form = UploadedMaterialForm(request.POST, request.FILES)
#         if form.is_valid():
#             material = form.save(commit=False)
#             material.course = course
#             material.save()
#             return redirect('course_detail', course_id=course.id)
#     else:
#         form = UploadedMaterialForm()
#     return render(request, 'courses/upload_material.html', {'form': form, 'course': course})

def extract_text_from_pdf(uploaded_file):
    text = ""
    reader = PdfReader(uploaded_file)
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(uploaded_file):
    text = ""
    document = Document(uploaded_file)
    for para in document.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_pptx(uploaded_file):
    text = ""
    presentation = Presentation(uploaded_file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

@login_required
def upload_material(request, course_id):
    course = get_object_or_404(Course, id=course_id, professor=request.user)
    
    if request.method == 'POST':
        form = UploadedMaterialForm(request.POST, request.FILES)
        
        if form.is_valid():
            material = form.save(commit=False)
            material.course = course

            # Process the uploaded file and convert it to a txt file
            uploaded_file = request.FILES['file']
            file_name, file_extension = os.path.splitext(uploaded_file.name)

            if file_extension == '.pdf':
                text = extract_text_from_pdf(uploaded_file)
            elif file_extension == '.docx':
                text = extract_text_from_docx(uploaded_file)
            elif file_extension == '.pptx':
                text = extract_text_from_pptx(uploaded_file)
            else:
                return HttpResponseBadRequest("Unsupported file format.")

            print("TEXT", text)
            material.text = text

            material.save()
            return redirect('course_detail', course_id=course.id)

    else:
        form = UploadedMaterialForm()

    return render(request, 'courses/upload_material.html', {'form': form, 'course': course})


@login_required
def material_detail(request, course_id, material_id):
    material = get_object_or_404(UploadedMaterial, id=material_id, course_id=course_id)
    return render(request, 'courses/material_detail.html', {'material': material})


@login_required
def delete_material(request, course_id, material_id):
    material = get_object_or_404(UploadedMaterial, id=material_id, course_id=course_id)

    # Check if the user is the professor of the course
    if material.course.professor != request.user:
        raise PermissionDenied

    if request.method == 'POST':
        # Delete the file associated with the material
        material.file.delete(save=False)  # save=False prevents the model from being saved after file deletion
        # Delete the material record
        material.delete()
        return redirect('course_detail', course_id=course_id)
    else:
        return HttpResponseForbidden('Only POST method is accepted')
